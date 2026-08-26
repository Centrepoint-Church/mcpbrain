"""Binary file text extractors for Drive sync.

Handles: PDF (text layer + optional OCR), DOCX (paragraphs + tables), XLSX/XLS
(structured tables, unbounded rows within a char budget — see tabular.py),
PPTX (slide text + tables), EML (headers + body as prose). All imports are
lazy so the module loads even if a dependency is missing — failed extractions
return "" (or [] for the tabular extractors), and now always log a warning
naming the extractor and the failure (B7): a corrupt file and an unsupported
one used to be indistinguishable, both silently yielding nothing.

Scanned/image-only PDFs are OCR'd page-by-page via the standalone `tesseract`
CLI (render page → PNG → `tesseract … stdout`). This deliberately does NOT use
pymupdf's built-in get_textpage_ocr, which needs MuPDF compiled with Tesseract
integration + TESSDATA_PREFIX — the pip wheel usually isn't, so it fails on a
plain `brew install tesseract`. The CLI path works with any tesseract install.
tesseract is an optional external binary (NOT a pip dependency); if it is absent,
image-only PDFs degrade gracefully to whatever text layer exists ('' when none),
but that degradation is now logged (A5) instead of silent.
"""

import io
import logging
import os
import shutil
import subprocess
import tempfile

from mcpbrain.sync.tabular import Table, normalise_rows

log = logging.getLogger(__name__)

_OCR_MIN_PAGE_CHARS = 20   # a page with fewer real chars is treated as image-only
_OCR_DPI = 200             # render resolution for OCR (quality vs speed)


# ---------------------------------------------------------------------------
# Partial-result signalling (I9)
# ---------------------------------------------------------------------------

class PartialTables(list):
    """`list[Table]` from an extraction that FAILED PARTWAY through the document.

    The multi-sheet/multi-slide extractors below deliberately keep whatever they
    had accumulated when an exception hit mid-iteration — better than nothing.
    But drive.upsert_file_chunks reads a SHORT chunk list as evidence that the
    document SHRANK and deletes the "orphaned" higher-index chunks (B5). So a
    transient failure on sheet 3 of 5 permanently deleted sheets 3-5's
    previously-good chunks, and nothing re-triggers extraction for a file whose
    metadata never changes again — a logged warning turned into irreversible
    content loss.

    A list/str SUBCLASS rather than a tuple or a wrapper dataclass: it compares
    and behaves exactly like the plain result, so every existing caller and test
    is untouched, and the one caller that needs to know asks `is_partial()`.
    """


class PartialText(str):
    """As PartialTables, for the text-returning extractors."""


def is_partial(result) -> bool:
    """True when `result` came from an extraction that died partway (see
    PartialTables). False for a complete result, and for None."""
    return isinstance(result, (PartialTables, PartialText))


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _is_scanned_pages(pages: list[str], chars_per_page_threshold: int = 50) -> bool:
    """The scanned/image-only decision, from already-extracted page text.

    Split out so the decision is computed ONCE per document (I7): extract_text_from_pdf
    already has every page's text and used to call is_scanned_pdf, which opened a
    SECOND fitz document and re-extracted all of it — doubling the cost of the
    corpus's most expensive extractor on every PDF, forever.
    """
    if not pages:
        return False
    total = sum(len(p or "") for p in pages)
    return (total / len(pages)) < chars_per_page_threshold


def is_scanned_pdf(content_bytes: bytes, *, chars_per_page_threshold: int = 50,
                   pages: list[str] | None = None) -> bool:
    """True when a PDF looks scanned/image-only (avg text-layer chars/page is low).

    Mirrors ops-brain's pdf_scanned_check. Used to decide whether OCR is worth
    attempting; returns False on any open error (caller falls back to text layer).

    `pages`, when given, is the document's already-extracted per-page text: the
    decision is made from it directly and `content_bytes` is never opened (I7).
    The DECISION is identical either way — both paths reduce to
    _is_scanned_pages over the same page text.
    """
    if pages is not None:
        return _is_scanned_pages(pages, chars_per_page_threshold)
    try:
        import fitz  # pymupdf
        doc = fitz.open(stream=content_bytes, filetype="pdf")
    except Exception as exc:
        log.debug("is_scanned_pdf: open failed: %s", exc)
        return False
    try:
        return _is_scanned_pages([page.get_text() or "" for page in doc],
                                 chars_per_page_threshold)
    except Exception as exc:
        log.debug("is_scanned_pdf: detection failed: %s", exc)
        return False
    finally:
        doc.close()


def extract_text_from_pdf(content_bytes: bytes) -> str:
    """PDF text via pymupdf; per-page OCR fallback (tesseract CLI) for scanned pages.

    Every path that yields less than the document contains now says so (A5).
    Previously a scanned PDF with tesseract absent returned '' in silence, and a
    per-page OCR timeout (120 s) fell back to an empty page_text unlogged — so
    the file simply had no chunks and nothing recorded why.

    `is_scanned_pdf` is now the single gate. It was dead code sitting beside a
    second, DIFFERENT inline heuristic (avg < 50 chars/page here, avg < 20
    inline), and two heuristics that can disagree — one of them unreachable —
    is worse than either alone. The new gate is slightly more willing to attempt
    OCR, which is the intended direction: a page with 30 characters is a scan
    with a caption.
    """
    try:
        import fitz  # pymupdf
        doc = fitz.open(stream=content_bytes, filetype="pdf")
    except Exception as exc:
        log.warning("pdf: open failed: %s", exc)
        return ""
    try:
        pages = [page.get_text() for page in doc]
        # `pages=` so the gate reuses the text we just extracted instead of
        # re-opening and re-parsing the whole document (I7).
        if not is_scanned_pdf(content_bytes, pages=pages):
            return "\n\n".join(pages)
        if not _tesseract_available():
            log.warning("pdf: looks scanned (%d pages, %d text chars) and "
                        "tesseract is unavailable — returning the text layer only",
                        len(pages), sum(len(p or "") for p in pages))
            return "\n\n".join(pages)
        out, ocr_failures = [], 0
        for i, page in enumerate(doc):
            page_text = (pages[i] if i < len(pages) else page.get_text()).strip()
            if len(page_text) >= _OCR_MIN_PAGE_CHARS:
                out.append(page_text)
                continue
            ocr = _ocr_page(page)
            if not ocr:
                ocr_failures += 1
            out.append(ocr or page_text)
        if ocr_failures:
            log.warning("pdf: OCR produced nothing for %d of %d pages "
                        "(timeout or render failure) — those pages are empty",
                        ocr_failures, len(out))
        return "\n\n".join(out)
    except Exception as exc:
        log.warning("pdf: extraction failed: %s", exc)
        return ""
    finally:
        doc.close()  # guaranteed close on every path once the doc is open


def _ocr_page(page) -> str:
    """Render one pymupdf page to PNG and OCR it via the tesseract CLI.

    Returns the recognised text, or '' on any failure (missing binary, render
    error, non-zero exit, timeout). Never raises — OCR is best-effort.
    """
    tess = _tesseract_bin()
    if not tess:
        return ""
    try:
        png = page.get_pixmap(dpi=_OCR_DPI).tobytes("png")
    except Exception as exc:
        log.debug("ocr: page render failed: %s", exc)
        return ""
    tmp = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(png)
            tmp = f.name
        proc = subprocess.run(
            [tess, tmp, "stdout", "-l", "eng"],
            capture_output=True, text=True, timeout=120,
        )
        return (proc.stdout or "").strip() if proc.returncode == 0 else ""
    except Exception as exc:
        log.debug("ocr: tesseract invocation failed: %s", exc)
        return ""
    finally:
        if tmp:
            try:
                os.unlink(tmp)
            except OSError:
                pass


# Common absolute locations checked in addition to PATH. The daemon runs under
# launchd/systemd with a MINIMAL PATH that usually excludes Homebrew, so
# shutil.which("tesseract") alone would miss a `brew install`ed binary.
_TESSERACT_FALLBACK_PATHS = (
    "/opt/homebrew/bin/tesseract",   # Apple Silicon Homebrew
    "/usr/local/bin/tesseract",      # Intel Homebrew
    "/usr/bin/tesseract",            # Linux distro packages
)

_tesseract_cache = None  # cached resolved path (str) or "" when absent


def _tesseract_bin() -> str:
    """Resolve the tesseract binary: PATH first, then known install locations.

    Returns the path, or '' if not found. Cached. The fallback paths matter
    because the daemon's launchd/systemd PATH typically omits Homebrew dirs.
    Set TESSERACT_BIN to override explicitly.
    """
    global _tesseract_cache
    if _tesseract_cache is None:
        env = os.environ.get("TESSERACT_BIN", "")
        found = env or shutil.which("tesseract") or ""
        if not found:
            for p in _TESSERACT_FALLBACK_PATHS:
                if os.path.exists(p):
                    found = p
                    break
        _tesseract_cache = found
        # Return the local `found`, not a re-read of the global: this now runs
        # inside the long-lived daemon, where ocr.tesseract_available() busts
        # the cache (sets it to None) from a background thread to force a
        # fresh resolve. A concurrent extraction thread's own call here could
        # land between the assignment above and a `return _tesseract_cache`
        # re-read, observing a bust from that other caller and returning None
        # even though THIS call just successfully resolved a real path.
        return found
    return _tesseract_cache


def _tesseract_available() -> bool:
    return bool(_tesseract_bin())


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def extract_text_from_docx(content_bytes: bytes) -> str:
    """Extract text from DOCX bytes, including table content."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content_bytes))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as exc:
        log.warning("docx: extraction failed: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# XLSX / XLS (structured tables — see mcpbrain.sync.tabular)
# ---------------------------------------------------------------------------

def _tables_from_grid(name: str, raw: list[list[str]],
                      char_budget: int) -> list[Table]:
    """Normalise one sheet's raw grid into at most one `Table`.

    Shared by the .xlsx and .xls readers so the two formats cannot drift on
    empty-row handling, column trimming or the budget rule.

    `char_budget` bounds the sheet over NON-EMPTY rows and replaces the flat
    200-rows-per-sheet cap (B1): 338 live files hit that cap, including several
    budgets and a general ledger, losing everything past row 200 per sheet while
    the same files bloated the store with empty cells.
    """
    rows = normalise_rows(raw)
    if not rows:
        return []
    header, data = rows[0], rows[1:]
    kept, used = [], 0
    for row in data:
        size = sum(len(c) + 1 for c in row)
        if used + size > char_budget and kept:
            break
        kept.append(row)
        used += size
    truncated = len(kept) < len(data)
    if truncated:
        log.warning("sheet %r truncated at %d of %d rows (char budget %d)",
                    name, len(kept), len(data), char_budget)
    return [Table(sheet=name, header=header, rows=kept,
                  rows_total=len(data), truncated=truncated)]


def extract_tables_from_xlsx(content_bytes: bytes, *, char_budget: int) -> list[Table]:
    """Read every sheet as a `Table`.

    Replaces extract_text_from_xlsx. Deliberately does NOT render: chunk
    boundaries are decided by tabular.render_chunks, so each chunk can repeat
    the header rather than orphaning it in chunk 0 (B2).

    A mid-iteration failure returns the sheets read so far as `PartialTables`
    (I9) so the caller does not mistake the short result for a shrunk document.
    """
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content_bytes), read_only=True,
                                    data_only=True)
    except Exception as exc:
        log.warning("xlsx: workbook open failed: %s", exc)
        return []
    tables: list[Table] = []
    partial = False
    try:
        for name in wb.sheetnames:
            raw = [[str(c) if c is not None else "" for c in row]
                   for row in wb[name].iter_rows(values_only=True)]
            tables.extend(_tables_from_grid(name, raw, char_budget))
    except Exception as exc:
        log.warning("xlsx: extraction failed after %d sheets: %s", len(tables), exc)
        partial = True
    finally:
        wb.close()
    return PartialTables(tables) if partial else tables


def _xls_cell_to_str(cell, datemode: int) -> str:
    """Render one xlrd cell as a string, honouring its Excel type.

    `sheet.row_values(r)` (the naive approach) returns raw Python floats for
    both numbers AND dates — a date comes back as an Excel serial number like
    45352.0, and a whole-number amount comes back as 500.0 with a spurious
    trailing '.0'. Both are real defects for "budgets, ledgers and risk
    assessments" (the stated reason this format was added at all), not
    cosmetic: a date column full of 5-digit serials is unusable, and every
    integral amount picks up a fake decimal. `cell.ctype` distinguishes them.
    """
    import xlrd

    if cell.ctype == xlrd.XL_CELL_EMPTY or cell.value in (None, ""):
        return ""
    if cell.ctype == xlrd.XL_CELL_DATE:
        try:
            dt = xlrd.xldate_as_datetime(cell.value, datemode)
        except (xlrd.XLDateError, ValueError, OverflowError) as exc:
            log.debug("xls: bad date serial %r: %s", cell.value, exc)
            return str(cell.value)
        if (dt.hour, dt.minute, dt.second, dt.microsecond) == (0, 0, 0, 0):
            return dt.date().isoformat()
        return dt.isoformat(sep=" ")
    if cell.ctype == xlrd.XL_CELL_NUMBER:
        v = cell.value
        return str(int(v)) if v == int(v) else str(v)
    if cell.ctype == xlrd.XL_CELL_BOOLEAN:
        return "TRUE" if cell.value else "FALSE"
    return str(cell.value)


def extract_tables_from_xls(content_bytes: bytes, *, char_budget: int) -> list[Table]:
    """Legacy .xls via xlrd, yielding the same `Table` shape as .xlsx.

    A2: .xls was dropped entirely. Declining a spreadsheet format is not
    defensible after B1 — budgets, ledgers and risk assessments are the
    highest-value tabular content in the corpus — and xlrd 2.0 exists precisely
    for this format (it dropped .xlsx, which openpyxl handles).
    """
    try:
        import xlrd
        book = xlrd.open_workbook(file_contents=content_bytes)
    except Exception as exc:
        log.warning("xls: workbook open failed: %s", exc)
        return []
    tables: list[Table] = []
    try:
        for sheet in book.sheets():
            raw = [[_xls_cell_to_str(sheet.cell(r, c), book.datemode)
                    for c in range(sheet.ncols)]
                   for r in range(sheet.nrows)]
            tables.extend(_tables_from_grid(sheet.name, raw, char_budget))
    except Exception as exc:
        log.warning("xls: extraction failed after %d sheets: %s", len(tables), exc)
        return PartialTables(tables)   # I9 — see PartialTables
    return tables


def extract_text_from_eml(content_bytes: bytes) -> str:
    """A .eml file's headers and body as prose.

    A2: .eml files in Drive were dropped. Stdlib `email` parses them with no new
    dependency. Deliberately extracted as a PROSE DOCUMENT rather than turned
    into a synthetic Gmail thread: a Drive file is not a mailbox message, and
    minting message_id/thread_id for it would put a second, unauthoritative
    identity into the same namespace the real Gmail sync owns.
    """
    try:
        from email import policy
        from email.parser import BytesParser
        msg = BytesParser(policy=policy.default).parsebytes(content_bytes)
    except Exception as exc:
        log.warning("eml: parse failed: %s", exc)
        return ""
    try:
        head = [f"{h}: {msg.get(h, '')}" for h in ("From", "To", "Cc", "Subject", "Date")
                if msg.get(h)]
        body = msg.get_body(preferencelist=("plain", "html"))
        text = body.get_content() if body is not None else ""
        if body is not None and body.get_content_type() == "text/html":
            from mcpbrain.sync.normalise import strip_html
            text = strip_html(text)
        return "\n".join(head) + "\n\n" + (text or "").strip()
    except Exception as exc:
        log.warning("eml: extraction failed: %s", exc)
        return ""


def extract_text_from_pptx(content_bytes: bytes) -> str:
    """Extract slide text from PPTX bytes: titles, body frames and table cells.

    Slides are separated by a labelled heading so chunk_text's paragraph split
    keeps slide boundaries, and so a recalled chunk says which slide it is from.
    """
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content_bytes))
    except Exception as exc:
        log.warning("pptx: presentation open failed: %s", exc)
        return ""
    parts: list[str] = []
    try:
        for n, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            lines.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            lines.append(" | ".join(cells))
            if lines:
                parts.append(f"Slide {n}\n" + "\n".join(lines))
    except Exception as exc:
        log.warning("pptx: extraction failed after %d slides: %s", len(parts), exc)
        return PartialText("\n\n".join(parts))   # I9 — see PartialTables
    return "\n\n".join(parts)
