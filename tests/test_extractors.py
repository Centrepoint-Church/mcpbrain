"""Tests for mcpbrain.sync.extractors — real in-memory files, no network."""

import io



# ---------------------------------------------------------------------------
# Helpers to build real binary files in memory
# ---------------------------------------------------------------------------

def _make_docx_bytes() -> bytes:
    from docx import Document
    doc = Document()
    doc.add_paragraph("Quarterly budget review")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Revenue"
    table.rows[0].cells[1].text = "Expenses"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes() -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Category", "Amount"])
    ws.append(["Salaries", 120000])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pdf_with_text_bytes() -> bytes:
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Budget report Q3")
    data = doc.tobytes()
    doc.close()
    return data


def _make_pdf_no_text_bytes() -> bytes:
    """A PDF with a page that has no text layer."""
    import fitz
    doc = fitz.open()
    doc.new_page()   # blank page — no text inserted
    data = doc.tobytes()
    doc.close()
    return data


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

def test_docx_roundtrip():
    """DOCX with a paragraph and a table row extracts both."""
    from mcpbrain.sync.extractors import extract_text_from_docx
    text = extract_text_from_docx(_make_docx_bytes())
    assert "Quarterly budget review" in text
    assert "Revenue" in text
    assert "Expenses" in text


def test_xlsx_roundtrip():
    """extract_text_from_xlsx no longer exists: chunk boundaries used to be
    decided here (a single markdown-rendered string per sheet), which orphaned
    the header in chunk 0 once chunk_text split it (B2). Chunking is now the
    renderer's job (tabular.render_chunks), so the extractor's contract is a
    structured Table, not text — assert the same content via the new shape."""
    from mcpbrain.sync.extractors import extract_tables_from_xlsx
    tables = extract_tables_from_xlsx(_make_xlsx_bytes(), char_budget=1_000_000)
    assert len(tables) == 1
    assert tables[0].sheet == "Budget"
    assert tables[0].header == ["Category", "Amount"]
    assert tables[0].rows == [["Salaries", "120000"]]


def test_pdf_text_layer():
    """A PDF with a real text layer extracts text without needing OCR."""
    from mcpbrain.sync.extractors import extract_text_from_pdf
    text = extract_text_from_pdf(_make_pdf_with_text_bytes())
    assert "Budget report Q3" in text


def test_pdf_no_text_layer_degrades_without_tesseract(monkeypatch):
    """Scanned/image PDF with no text layer: when tesseract unavailable, returns
    empty string (or whitespace) without raising. Documents graceful degradation."""
    import mcpbrain.sync.extractors as extractors_mod
    monkeypatch.setattr(extractors_mod, "_tesseract_cache", False)

    text = extractors_mod.extract_text_from_pdf(_make_pdf_no_text_bytes())

    # Must not raise; result is empty (no text layer, no OCR)
    assert isinstance(text, str)
    assert text.strip() == ""


def test_extractors_return_empty_on_garbage():
    """Garbage bytes fed to each extractor returns '' (or [] for the tabular
    extractors, which return a list of Table, not text) without crashing.
    extract_text_from_xlsx no longer exists (see test_xlsx_roundtrip); the
    replacement extract_tables_from_xlsx keeps the same graceful-failure
    contract, just with a list return type."""
    from mcpbrain.sync.extractors import (
        extract_text_from_pdf,
        extract_text_from_docx,
        extract_tables_from_xlsx,
    )
    garbage = b"not a real file"
    assert extract_text_from_pdf(garbage) == ""
    assert extract_text_from_docx(garbage) == ""
    assert extract_tables_from_xlsx(garbage, char_budget=1_000_000) == []


# ---------------------------------------------------------------------------
# Scanned-PDF detection + tesseract OCR (Q5)
# ---------------------------------------------------------------------------

def _make_pdf_long_text_bytes() -> bytes:
    import fitz
    doc = fitz.open(); page = doc.new_page()
    page.insert_text((72, 72),
                     "This is a budget report with plenty of real text on the page, "
                     "well above the scanned-PDF character threshold per page.")
    data = doc.tobytes(); doc.close()
    return data


def test_is_scanned_pdf_true_for_blank():
    from mcpbrain.sync.extractors import is_scanned_pdf
    assert is_scanned_pdf(_make_pdf_no_text_bytes()) is True


def test_is_scanned_pdf_false_for_text_pdf():
    from mcpbrain.sync.extractors import is_scanned_pdf
    assert is_scanned_pdf(_make_pdf_long_text_bytes()) is False


def test_digital_pdf_returns_text_layer():
    from mcpbrain.sync.extractors import extract_text_from_pdf
    out = extract_text_from_pdf(_make_pdf_long_text_bytes())
    assert "budget report" in out.lower()


def test_scanned_pdf_degrades_to_empty_without_tesseract(monkeypatch):
    """No tesseract resolvable → a scanned PDF yields '' (graceful), never raises."""
    import mcpbrain.sync.extractors as ex
    monkeypatch.setattr(ex, "_tesseract_bin", lambda: "")
    out = ex.extract_text_from_pdf(_make_pdf_no_text_bytes())
    assert out.strip() == ""


def test_scanned_pdf_uses_ocr_output_when_available(monkeypatch):
    """Control-flow: a no-text page routes through _ocr_page and its text is used."""
    import mcpbrain.sync.extractors as ex
    monkeypatch.setattr(ex, "_tesseract_cache", True)          # pretend tesseract present
    monkeypatch.setattr(ex, "_ocr_page", lambda page: "OCR RECOVERED BUDGET")
    out = ex.extract_text_from_pdf(_make_pdf_no_text_bytes())
    assert "OCR RECOVERED BUDGET" in out


def test_ocr_roundtrip_with_real_tesseract():
    """End-to-end OCR of an image-only PDF — runs only where tesseract is installed."""
    import shutil
    import pytest
    if not shutil.which("tesseract"):
        pytest.skip("tesseract not installed")
    from PIL import Image, ImageDraw, ImageFont
    import fitz
    from mcpbrain.sync.extractors import extract_text_from_pdf, is_scanned_pdf

    img = Image.new("RGB", (900, 240), "white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 96)
    except Exception:
        font = ImageFont.load_default()
    draw.text((30, 60), "CENTREPOINT", fill="black", font=font)
    pbuf = io.BytesIO(); img.save(pbuf, format="PNG")

    doc = fitz.open(); page = doc.new_page(width=900, height=240)
    page.insert_image(fitz.Rect(0, 0, 900, 240), stream=pbuf.getvalue())
    pdf = doc.tobytes(); doc.close()

    assert is_scanned_pdf(pdf) is True
    out = extract_text_from_pdf(pdf)
    assert "centrepoint" in out.lower()


# ---------------------------------------------------------------------------
# Structured tables (A2/B1/B2), PPTX, legacy XLS, EML (A2), visible failures (B7)
# ---------------------------------------------------------------------------

def test_xlsx_yields_structured_tables():
    """extract_text_from_xlsx is replaced by extract_tables_from_xlsx: chunk
    boundaries are decided later, by the renderer, so each chunk can repeat the
    header instead of orphaning it in chunk 0 (B2)."""
    import io

    import openpyxl

    from mcpbrain.sync.extractors import extract_tables_from_xlsx

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Item", "Amount"])
    ws.append(["Rent", 500])
    ws.append([None, None])          # empty row — must be dropped
    ws.append(["Power", 120])
    buf = io.BytesIO()
    wb.save(buf)

    tables = extract_tables_from_xlsx(buf.getvalue(), char_budget=1_000_000)

    assert len(tables) == 1
    assert tables[0].sheet == "Budget"
    assert tables[0].header == ["Item", "Amount"]
    assert tables[0].rows == [["Rent", "500"], ["Power", "120"]]
    assert tables[0].truncated is False


def test_xlsx_keeps_rows_past_the_old_200_row_cap():
    """338 live files hit the old cap — budgets, a general ledger, risk
    assessments — losing every row past 200 per sheet."""
    import io

    import openpyxl

    from mcpbrain.sync.extractors import extract_tables_from_xlsx

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Item", "Amount"])
    for i in range(400):
        ws.append([f"Item {i}", i])
    buf = io.BytesIO()
    wb.save(buf)

    tables = extract_tables_from_xlsx(buf.getvalue(), char_budget=1_000_000)

    assert len(tables[0].rows) == 400
    assert tables[0].truncated is False


def test_pptx_text_is_extracted():
    """A2: presentationml.presentation is advertised in _MIME_EXTRACTION_META but
    reachable by no fetcher, so every .pptx was dropped. Verified live: 0 chunks
    for .pptx against 28 for native Google Slides."""
    import io

    from pptx import Presentation

    from mcpbrain.sync.extractors import extract_text_from_pptx

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Q3 Ministry Review"
    buf = io.BytesIO()
    prs.save(buf)

    assert "Q3 Ministry Review" in extract_text_from_pptx(buf.getvalue())


def test_legacy_xls_yields_the_same_table_shape_as_xlsx():
    """A2: .xls was dropped entirely. Declining a SPREADSHEET format would be
    indefensible after B1 established that budgets and ledgers are the
    highest-value tabular content, and xlrd 2.0 exists purely to read .xls.

    The fixture is a small binary .xls checked into tests/fixtures/, generated
    once via `uv run --with xlwt python -c "..."` (xlwt is a dev-only,
    write-only helper — not a project dependency — so it is never imported
    here; only xlrd, already in pyproject.toml, is used at test time).
    """
    from pathlib import Path

    from mcpbrain.sync.extractors import extract_tables_from_xls

    raw = (Path(__file__).parent / "fixtures" / "legacy_budget.xls").read_bytes()

    tables = extract_tables_from_xls(raw, char_budget=1_000_000)

    assert len(tables) == 1
    assert tables[0].sheet == "Budget"
    assert tables[0].header == ["Item", "Amount"]
    assert tables[0].rows == [["Rent", "500"]]


def test_legacy_xls_renders_real_numbers_and_dates_correctly():
    """Important #3 (review): sheet.row_values() — the naive xlrd read — returns
    raw Python floats for BOTH numbers and dates, so a real numeric/date .xls
    (as opposed to the all-strings fixture above, which is exactly why this
    slipped through) rendered as [['Item', 'Amount', 'Date'],
    ['Rent', '500.0', '45352.0']] instead of a clean integer and an ISO date —
    corrupting the "budgets, ledgers and risk assessments" this format exists
    for. The fixture has a real numeric int-valued cell (500), a real
    non-integral numeric cell (42.5) and a real Excel date cell (2024-03-01,
    written with a date number format so xlrd types it XL_CELL_DATE), built
    the same way as legacy_budget.xls (uv run --with xlwt ...).
    """
    from pathlib import Path

    from mcpbrain.sync.extractors import extract_tables_from_xls

    raw = (Path(__file__).parent / "fixtures" / "legacy_budget_typed.xls").read_bytes()

    tables = extract_tables_from_xls(raw, char_budget=1_000_000)

    assert len(tables) == 1
    assert tables[0].header == ["Item", "Amount", "Date"]
    rent, bus = tables[0].rows
    assert rent[:2] == ["Rent", "500"], "integral amount must not read '500.0'"
    assert rent[2] == "2024-03-01", "a date cell must render as a date, not an Excel serial"
    assert bus[:2] == ["Bus fare", "42.5"], "a non-integral amount must round-trip cleanly"


def test_eml_is_extracted_as_prose_with_its_headers():
    """A2: .eml files in Drive were dropped. Stdlib `email` parses them — zero
    new dependencies — and they become prose documents, NOT synthetic Gmail
    threads (that would be scope creep into the sync layer's identity model)."""
    from mcpbrain.sync.extractors import extract_text_from_eml

    raw = (b"From: sam@example.com\r\n"
           b"To: josh@centrepoint.church\r\n"
           b"Subject: Hall B booking\r\n"
           b"Date: Tue, 02 Jun 2026 16:30:01 +0800\r\n"
           b"Content-Type: text/plain; charset=utf-8\r\n\r\n"
           b"Confirmed for Sunday the 8th.\r\n")

    text = extract_text_from_eml(raw)

    assert "Subject: Hall B booking" in text
    assert "sam@example.com" in text
    assert "Confirmed for Sunday the 8th." in text


def test_a_multipart_eml_prefers_the_plain_text_part():
    from mcpbrain.sync.extractors import extract_text_from_eml

    raw = (b"Subject: Multi\r\n"
           b'Content-Type: multipart/alternative; boundary="b"\r\n\r\n'
           b"--b\r\nContent-Type: text/plain\r\n\r\nplain body here\r\n"
           b"--b\r\nContent-Type: text/html\r\n\r\n<p>html body</p>\r\n--b--\r\n")

    text = extract_text_from_eml(raw)

    assert "plain body here" in text
    assert "<p>" not in text


def test_pptx_extraction_failure_returns_empty_and_logs(caplog):
    from mcpbrain.sync.extractors import extract_text_from_pptx

    with caplog.at_level("WARNING"):
        assert extract_text_from_pptx(b"not a pptx") == ""

    assert any("pptx" in r.message for r in caplog.records), (
        "a failed extraction must leave a trace; eight sites in this module "
        "returned '' with no log line at all (B7)"
    )


def test_a_scanned_pdf_with_no_ocr_available_is_reported_not_silently_empty(monkeypatch, caplog):
    """A5: with tesseract absent a fully-scanned PDF returns '' with no warning
    at all — no chunks, no log line, nothing to explain the absence."""
    import fitz

    from mcpbrain.sync import extractors

    monkeypatch.setattr(extractors, "_tesseract_available", lambda: False)
    doc = fitz.open()
    doc.new_page()
    data = doc.tobytes()

    with caplog.at_level("WARNING"):
        text = extractors.extract_text_from_pdf(data)

    assert text.strip() == ""
    assert any("scanned" in r.message.lower() and "tesseract" in r.message.lower()
               for r in caplog.records), (
        f"no warning explained the empty result: {[r.message for r in caplog.records]}"
    )


def test_a_failed_ocr_page_is_logged(monkeypatch, caplog):
    """A5: per-page OCR failure/timeout returns '' and falls back to page_text,
    so a timed-out page yields nothing, unlogged."""
    import fitz

    from mcpbrain.sync import extractors

    monkeypatch.setattr(extractors, "_tesseract_available", lambda: True)
    monkeypatch.setattr(extractors, "_ocr_page", lambda page: "")
    doc = fitz.open()
    doc.new_page()

    with caplog.at_level("WARNING"):
        extractors.extract_text_from_pdf(doc.tobytes())

    assert any("ocr" in r.message.lower() for r in caplog.records)


def test_is_scanned_pdf_is_either_used_or_gone():
    """A5: is_scanned_pdf is defined but never called; the real gate is an
    inline char-count heuristic. Two heuristics that can disagree, one of them
    dead, is worse than either alone."""
    import subprocess

    out = subprocess.run(["git", "grep", "-n", "is_scanned_pdf", "--", "mcpbrain/"],
                         capture_output=True, text=True).stdout
    call_sites = [ln for ln in out.splitlines() if "def is_scanned_pdf" not in ln]

    assert call_sites, "is_scanned_pdf is still dead code"


def test_a_pdf_is_parsed_once_not_twice():
    """I7: extract_text_from_pdf extracted `pages` itself and then called
    is_scanned_pdf(content_bytes), which opened a SECOND fitz document and
    re-extracted every page just to make the scanned/not-scanned decision —
    doubling the cost of the corpus's most expensive extractor on every PDF,
    forever. The decision is unchanged; it is now computed from the pages we
    already have."""
    import fitz

    from mcpbrain.sync.extractors import extract_text_from_pdf

    # Built BEFORE the patch — the fixture helper opens fitz itself to make the
    # PDF, which would otherwise be counted.
    pdf = _make_pdf_long_text_bytes()
    opens = []
    real_open = fitz.open

    def counting_open(*a, **kw):
        opens.append(1)
        return real_open(*a, **kw)

    fitz.open = counting_open
    try:
        text = extract_text_from_pdf(pdf)
    finally:
        fitz.open = real_open

    assert "Hello" in text or text.strip(), "extraction still has to work"
    assert len(opens) == 1, f"the document was opened {len(opens)} times"


def test_the_scanned_decision_is_identical_whichever_way_it_is_computed():
    """The refactor must not move the gate: is_scanned_pdf(bytes) and
    is_scanned_pdf(bytes, pages=…) must agree for both a text PDF and a blank one."""
    import fitz

    from mcpbrain.sync.extractors import is_scanned_pdf

    for pdf, expected in ((_make_pdf_long_text_bytes(), False),
                          (_make_pdf_no_text_bytes(), True)):
        doc = fitz.open(stream=pdf, filetype="pdf")
        pages = [p.get_text() for p in doc]
        doc.close()
        assert is_scanned_pdf(pdf) is expected
        assert is_scanned_pdf(pdf, pages=pages) is expected


# ---------------------------------------------------------------------------
# I9: partial-result signalling
# ---------------------------------------------------------------------------

def _make_multi_sheet_xlsx_bytes(sheets: int = 3) -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    for n in range(sheets):
        ws = wb.active if n == 0 else wb.create_sheet()
        ws.title = f"Sheet{n}"
        ws.append(["Item", "Amount"])
        ws.append([f"Row {n}", 100 + n])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_a_mid_workbook_failure_returns_a_partial_result(monkeypatch):
    """I9: the extractor keeps the sheets it got (better than nothing), but the
    CALLER has to be able to tell that apart from a genuinely short workbook —
    otherwise drive.upsert_file_chunks deletes the chunks the failed extraction
    never reached."""
    from mcpbrain.sync import extractors
    from mcpbrain.sync.extractors import extract_tables_from_xlsx, is_partial
    from mcpbrain.sync.tabular import Table

    calls = {"n": 0}

    def exploding(name, raw, char_budget):
        calls["n"] += 1
        if calls["n"] > 1:
            raise RuntimeError("simulated mid-workbook failure")
        return [Table(sheet=name, header=["a"], rows=[["1"]], rows_total=1)]

    monkeypatch.setattr(extractors, "_tables_from_grid", exploding)

    tables = extract_tables_from_xlsx(_make_multi_sheet_xlsx_bytes(),
                                      char_budget=1_000_000)

    assert len(tables) == 1, "the sheets read before the failure are kept"
    assert is_partial(tables) is True


def test_a_complete_workbook_is_not_marked_partial():
    from mcpbrain.sync.extractors import extract_tables_from_xlsx, is_partial

    tables = extract_tables_from_xlsx(_make_multi_sheet_xlsx_bytes(),
                                      char_budget=1_000_000)

    assert len(tables) >= 2
    assert is_partial(tables) is False


def test_a_partial_result_still_behaves_as_the_plain_list_it_wraps():
    """PartialTables/PartialText are list/str SUBCLASSES precisely so no existing
    caller or test has to change."""
    from mcpbrain.sync.extractors import PartialTables, PartialText, is_partial

    assert PartialTables([1, 2]) == [1, 2]
    assert PartialText("abc") == "abc"
    assert is_partial([1, 2]) is False
    assert is_partial("abc") is False
    assert is_partial(None) is False


def test_a_mid_deck_pptx_failure_returns_a_partial_string(monkeypatch):
    from mcpbrain.sync import extractors
    from mcpbrain.sync.extractors import extract_text_from_pptx, is_partial

    class _Run:
        text = "Slide body text"

    class _Para:
        runs = [_Run()]

    class _Frame:
        paragraphs = [_Para()]

    class _Shape:
        has_text_frame = True
        text_frame = _Frame()

    class _Slide:
        shapes = [_Shape()]

    class _Slides:
        def __iter__(self):
            yield _Slide()
            raise RuntimeError("simulated mid-deck failure")

    class _Prs:
        slides = _Slides()

    monkeypatch.setitem(__import__("sys").modules, "pptx",
                        type("m", (), {"Presentation": lambda _b: _Prs()}))

    out = extract_text_from_pptx(b"fake")

    assert "Slide body text" in out
    assert is_partial(out) is True
    assert extractors.is_partial(out) is True
